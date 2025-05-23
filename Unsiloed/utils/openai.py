import os
import base64
import json
from typing import List, Dict, Any, Optional
from openai import OpenAI
import logging
import concurrent.futures
import PyPDF2
from dotenv import load_dotenv
import numpy as np
import cv2
from functools import lru_cache
import hashlib
import threading
from collections import OrderedDict
import time
import re

load_dotenv()

logger = logging.getLogger(__name__)

# Configure logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)

# Instead of initializing at import, create a function to get the client
client = None
client_lock = threading.Lock()

def get_openai_client():
    """Get an OpenAI client with proper configuration and caching"""
    global client
    if client is not None:
        return client
        
    with client_lock:
        if client is not None:  # Double-check pattern
            return client
            
        try:
            api_key = os.environ.get("OPENAI_API_KEY")
            if not api_key:
                logger.error("OPENAI_API_KEY environment variable is not set")
                raise ValueError("OPENAI_API_KEY environment variable is not set")

            # Create client with optimized parameters
            client = OpenAI(
                api_key=api_key,
                timeout=30.0,  # Reduced timeout for faster failure detection
                max_retries=2,  # Reduced retries for faster fallback
            )
            
            # Test the client with a lightweight operation
            client.models.list(limit=1)  # Only fetch one model to test
            return client

        except Exception as e:
            logger.error(f"Error initializing OpenAI client: {str(e)}")
            return None

# Optimize cache with better memory management
class ThreadSafeLRUCache:
    def __init__(self, maxsize=100):
        self.cache = OrderedDict()
        self.maxsize = maxsize
        self.lock = threading.Lock()
        self.hits = 0
        self.misses = 0
    
    def get(self, key: str) -> Optional[str]:
        with self.lock:
            if key in self.cache:
                # Move to end (most recently used)
                value = self.cache.pop(key)
                self.cache[key] = value
                self.hits += 1
                return value
            self.misses += 1
            return None
    
    def put(self, key: str, value: str):
        with self.lock:
            if key in self.cache:
                self.cache.pop(key)
            elif len(self.cache) >= self.maxsize:
                # Remove least recently used item
                self.cache.popitem(last=False)
            self.cache[key] = value
    
    def clear(self):
        with self.lock:
            self.cache.clear()
            self.hits = 0
            self.misses = 0
    
    def get_stats(self):
        with self.lock:
            total = self.hits + self.misses
            hit_rate = (self.hits / total * 100) if total > 0 else 0
            return {
                "hits": self.hits,
                "misses": self.misses,
                "hit_rate": hit_rate,
                "size": len(self.cache)
            }

# Initialize thread-safe cache with monitoring
response_cache = ThreadSafeLRUCache(maxsize=100)

def get_cached_openai_response(text_hash: str, prompt: str) -> Optional[str]:
    """Get cached OpenAI response using thread-safe cache"""
    cache_key = f"{text_hash}:{prompt}"
    return response_cache.get(cache_key)

def cache_openai_response(text_hash: str, prompt: str, response: str):
    """Cache OpenAI response using thread-safe cache"""
    cache_key = f"{text_hash}:{prompt}"
    response_cache.put(cache_key, response)
    
    # Log cache stats periodically
    stats = response_cache.get_stats()
    if stats["hits"] % 100 == 0:  # Log every 100 hits
        logger.info(f"Cache stats: {stats}")

def get_text_hash(text: str) -> str:
    """Generate a hash for text to use as cache key"""
    return hashlib.md5(text.encode()).hexdigest()

def encode_image_to_base64(image_path):
    """
    Encode an image to base64.

    Args:
        image_path: Path to the image file or numpy array

    Returns:
        Base64 encoded string of the image
    """
    logger.debug("Encoding image to base64")

    # Handle numpy array (from CV2)
    if isinstance(image_path, np.ndarray):
        success, buffer = cv2.imencode(".jpg", image_path)
        if not success:
            raise ValueError("Failed to encode image")
        return base64.b64encode(buffer).decode("utf-8")

    # Handle file path
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def create_extraction_prompt(schema: Dict[str, Any], page_count: int) -> str:
    """
    Create a prompt instructing the model how to extract data according to the schema.

    Args:
        schema: JSON schema defining the structure
        page_count: Number of pages in the document

    Returns:
        Prompt string for the model
    """
    # Create a compact JSON representation of the schema
    schema_str = json.dumps(schema, indent=2)

    prompt = f"""
    You are an expert at extracting structured data from documents.
    
    I have a document with {page_count} pages that has been converted to images. I need you to extract specific information from these images according to the following JSON schema:
    
    {schema_str}
    
    Please follow these instructions carefully:
    
    1. Examine all {page_count} images thoroughly to find the requested information.
    2. Extract the exact text from the document that matches each field in the schema.
    3. If you cannot find information for a specific field in any of the pages, return an empty string or null value for that field.
    4. For array fields, include all instances found throughout the document.
    5. Maintain the structure defined in the schema exactly.
    6. Return only the extracted data as a valid JSON object, matching the structure of the schema.
    7. Do not add any explanatory text or notes outside the JSON structure.
    8. Be precise and accurate in your extraction.
    9. If text is unclear or ambiguous, make your best guess based on context.
    10. For dates, numbers, and other formatted data, maintain the format as shown in the document.
    11. IMPORTANT: Your response MUST be a valid JSON object that exactly matches the structure of the provided schema.
    12. IMPORTANT: Do not include any explanations, just return the JSON object.
    
    Your response should be a valid JSON object containing only the extracted data.
    """

    return prompt


def semantic_chunk_with_structured_output(text: str) -> List[Dict[str, Any]]:
    """
    Use OpenAI API to create semantic chunks from text.
    Optimized with caching, parallel processing, and retry logic.
    """
    # Generate hash for caching
    text_hash = get_text_hash(text)
    
    # Check cache first
    cached_response = get_cached_openai_response(text_hash, "semantic_chunk")
    if cached_response:
        return json.loads(cached_response)

    # If text is too long, split it first using a simpler method
    # and then process each part in parallel
    if len(text) > 25000:
        logger.info("Text too long for direct semantic chunking, applying parallel processing")
        return process_long_text_semantically(text)

    max_retries = 2
    retry_delay = 1.0
    
    for attempt in range(max_retries):
        try:
            # Get the OpenAI client
            openai_client = get_openai_client()
            if not openai_client:
                raise ValueError("Failed to initialize OpenAI client")

            # Split text into sections first
            sections = []
            current_section = []
            lines = text.split('. ')
            
            for line in lines:
                # Check if line starts with a numbered heading (e.g., "6.1", "6.2", etc.)
                if re.match(r'^\d+\.\d+\s+[A-Z]', line):
                    if current_section:
                        sections.append('. '.join(current_section) + '.')
                    current_section = [line]
                else:
                    current_section.append(line)
            
            # Add the last section
            if current_section:
                sections.append('. '.join(current_section) + '.')

            # Process each section separately
            all_chunks = []
            current_position = 0

            for section in sections:
                # Split section into paragraphs and bullet points
                paragraphs = []
                current_paragraph = []
                bullet_points = []
                current_bullet = []
                
                # First, identify section heading
                section_match = re.match(r'^(\d+\.\d+\s+[A-Z].*?)(?=\s|$)', section)
                section_heading = section_match.group(1) if section_match else ""
                
                # Process the rest of the section
                remaining_text = section[len(section_heading):].strip() if section_heading else section
                
                # Split into lines while preserving bullet points
                lines = []
                for line in remaining_text.split('. '):
                    if line.strip().startswith('●'):
                        if current_paragraph:
                            lines.append('. '.join(current_paragraph) + '.')
                            current_paragraph = []
                        lines.append(line)
                    else:
                        current_paragraph.append(line)
                        if len(' '.join(current_paragraph)) > 200:  # Split long paragraphs
                            lines.append('. '.join(current_paragraph) + '.')
                            current_paragraph = []
                
                # Add any remaining paragraph
                if current_paragraph:
                    lines.append('. '.join(current_paragraph) + '.')

                # Process lines into chunks
                current_chunk = []
                for line in lines:
                    if line.strip().startswith('●'):
                        # If we have a current chunk, add it to paragraphs
                        if current_chunk:
                            paragraphs.append(' '.join(current_chunk))
                            current_chunk = []
                        # Start a new bullet point chunk
                        current_chunk = [line]
                    else:
                        current_chunk.append(line)
                        if len(' '.join(current_chunk)) > 200:  # Split long chunks
                            paragraphs.append(' '.join(current_chunk))
                            current_chunk = []
                
                # Add any remaining chunk
                if current_chunk:
                    paragraphs.append(' '.join(current_chunk))

                # Add section heading to first chunk
                if section_heading:
                    if paragraphs:
                        paragraphs[0] = section_heading + ' ' + paragraphs[0]
                    else:
                        paragraphs.append(section_heading + '.')

                # Process each paragraph
                for chunk_text in paragraphs:
                    # Skip empty chunks or chunks that are too short
                    if not chunk_text or len(chunk_text) < 20:
                        continue
                    
                    # Ensure chunk ends with proper sentence boundary
                    if not any(chunk_text.rstrip().endswith(p) for p in ['.', '!', '?']):
                        continue
                        
                    # Find the chunk in the original text to get accurate character positions
                    start_position = text.find(chunk_text, current_position)
                    if start_position == -1:
                        # If exact match not found, use approximate position
                        start_position = current_position

                    end_position = start_position + len(chunk_text)

                    all_chunks.append(
                        {
                            "text": chunk_text,
                            "metadata": {
                                "title": f"Chunk {len(all_chunks) + 1}",
                                "position": "middle",
                                "start_char": start_position,
                                "end_char": end_position,
                                "strategy": "semantic",
                            },
                        }
                    )

                    current_position = end_position

            # Cache the response
            cache_openai_response(text_hash, "semantic_chunk", json.dumps({"chunks": all_chunks}))

            # Sort chunks by start position
            all_chunks.sort(key=lambda x: x["metadata"]["start_char"])

            # Merge small chunks with adjacent chunks, but only if they're related
            if len(all_chunks) > 1:
                merged_chunks = []
                current_chunk = all_chunks[0]
                
                for next_chunk in all_chunks[1:]:
                    # Only merge if chunks are small and related
                    if (len(current_chunk["text"]) < 100 and len(next_chunk["text"]) < 100 and
                        current_chunk["metadata"]["end_char"] == next_chunk["metadata"]["start_char"]):
                        # Merge chunks with proper spacing
                        current_chunk["text"] = current_chunk["text"].rstrip() + " " + next_chunk["text"].lstrip()
                        current_chunk["metadata"]["end_char"] = next_chunk["metadata"]["end_char"]
                    else:
                        merged_chunks.append(current_chunk)
                        current_chunk = next_chunk
                
                merged_chunks.append(current_chunk)
                all_chunks = merged_chunks

            return all_chunks

        except Exception as e:
            logger.error(f"Error in semantic chunking (attempt {attempt + 1}/{max_retries}): {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(retry_delay * (attempt + 1))  # Exponential backoff
                continue
            # Instead of falling back to paragraph chunking, raise the error
            raise ValueError(f"Failed to process text after {max_retries} attempts: {str(e)}")

def process_long_text_semantically(text: str) -> List[Dict[str, Any]]:
    """
    Process long text by splitting into sections and processing each section separately.
    Optimized for better performance with improved section management.
    """
    # Split text into sections based on numbered headings
    sections = []
    current_section = []
    lines = text.split('. ')
    
    for line in lines:
        # Check if line starts with a numbered heading (e.g., "6.1", "6.2", etc.)
        if re.match(r'^\d+\.\d+\s+[A-Z]', line):
            if current_section:
                sections.append('. '.join(current_section) + '.')
            current_section = [line]
        else:
            current_section.append(line)
    
    # Add the last section
    if current_section:
        sections.append('. '.join(current_section) + '.')
    
    # Process each section separately
    results = []
    current_position = 0
    
    for section in sections:
        # Process each section with semantic chunking
        section_chunks = semantic_chunk_with_structured_output(section)
        
        # Adjust positions to be relative to the full text
        for chunk in section_chunks:
            chunk["metadata"]["start_char"] += current_position
            chunk["metadata"]["end_char"] += current_position
            results.append(chunk)
        
        current_position += len(section)
    
    # Sort chunks by start position
    results.sort(key=lambda x: x["metadata"]["start_char"])
    
    # Merge overlapping chunks more efficiently, but never merge across section boundaries
    if not results:
        return []
        
    merged_chunks = [results[0]]
    
    for chunk in results[1:]:
        prev_chunk = merged_chunks[-1]
        
        # Check if chunks are from the same section
        prev_section = re.search(r'^\d+\.\d+\s+[A-Z]', prev_chunk["text"])
        curr_section = re.search(r'^\d+\.\d+\s+[A-Z]', chunk["text"])
        
        # Only merge if chunks are from the same section
        if prev_section and curr_section and prev_section.group() == curr_section.group():
            # Check for overlap
            if chunk["metadata"]["start_char"] <= prev_chunk["metadata"]["end_char"]:
                # Calculate overlap size
                overlap_size = prev_chunk["metadata"]["end_char"] - chunk["metadata"]["start_char"]
                
                # Only merge if overlap is significant (more than 50% of the smaller chunk)
                if overlap_size > min(len(prev_chunk["text"]), len(chunk["text"])) * 0.5:
                    # Merge chunks with proper spacing
                    prev_chunk["text"] = prev_chunk["text"].rstrip() + "\n" + chunk["text"].lstrip()
                    prev_chunk["metadata"]["end_char"] = chunk["metadata"]["end_char"]
                else:
                    merged_chunks.append(chunk)
            else:
                merged_chunks.append(chunk)
        else:
            # If chunks are from different sections, don't merge them
            merged_chunks.append(chunk)
    
    return merged_chunks


def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Extract text from PDF file with optimized parallel processing.
    Uses memory-efficient streaming and parallel page processing.
    """
    try:
        # Open PDF in binary mode for streaming
        with open(pdf_path, "rb") as file:
            # Create PDF reader with streaming mode
            reader = PyPDF2.PdfReader(file)
            total_pages = len(reader.pages)
            
            # Process pages in parallel with optimal worker count
            max_workers = min(total_pages, 4)  # Limit to 4 workers for optimal performance
            chunk_size = max(1, total_pages // max_workers)  # Calculate optimal chunk size
            
            # Use a list to store text chunks with pre-allocated size
            text_chunks = [None] * total_pages
            
            with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                def process_page_range(start_idx, end_idx):
                    try:
                        results = []
                        for i in range(start_idx, min(end_idx, total_pages)):
                            page = reader.pages[i]
                            text = page.extract_text()
                            if text and text.strip():  # Only store non-empty text
                                # Clean up the text
                                text = text.replace('\n', ' ')  # Replace newlines with spaces
                                text = ' '.join(text.split())  # Normalize whitespace
                                # Fix common sentence boundary issues
                                text = text.replace(' .', '.').replace(' ,', ',')
                                text = text.replace('  ', ' ')
                                # Ensure proper spacing after sentence endings
                                text = text.replace('.', '. ').replace('!', '! ').replace('?', '? ')
                                # Fix encoding issues
                                text = text.encode('ascii', 'ignore').decode('ascii')
                                # Fix bullet points
                                text = text.replace('•', '●').replace('·', '●')
                                text = ' '.join(text.split())  # Normalize again after fixes
                                results.append((i, text))
                        return results
                    except Exception as e:
                        logger.error(f"Error processing pages {start_idx}-{end_idx}: {str(e)}")
                        return []
                
                # Submit page ranges for parallel processing
                futures = []
                for i in range(0, total_pages, chunk_size):
                    end_idx = min(i + chunk_size, total_pages)
                    futures.append(executor.submit(process_page_range, i, end_idx))
                
                # Collect results in order
                for future in concurrent.futures.as_completed(futures):
                    for page_idx, text in future.result():
                        text_chunks[page_idx] = text
        
        # Join all text chunks with proper spacing, filtering out None values
        full_text = " ".join(chunk for chunk in text_chunks if chunk)
        
        # Final cleanup of the text
        full_text = ' '.join(full_text.split())  # Normalize whitespace
        full_text = full_text.replace(' .', '.').replace(' ,', ',')  # Fix common spacing issues
        full_text = full_text.replace('  ', ' ')  # Remove double spaces
        full_text = full_text.encode('ascii', 'ignore').decode('ascii')  # Fix encoding issues
        
        return full_text
        
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        raise


def extract_text_from_docx(docx_path: str) -> str:
    """
    Extract text from a DOCX file.

    Args:
        docx_path: Path to the DOCX file

    Returns:
        Extracted text from the DOCX
    """
    try:
        import docx  # python-docx package

        doc = docx.Document(docx_path)
        full_text = []

        # Extract text from paragraphs
        for para in doc.paragraphs:
            full_text.append(para.text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        raise


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extract text from a PPTX file.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        Extracted text from the PPTX
    """
    try:
        from pptx import Presentation  # python-pptx package

        presentation = Presentation(pptx_path)
        full_text = []

        # Loop through slides
        for slide_number, slide in enumerate(presentation.slides, 1):
            slide_text = []
            slide_text.append(f"Slide {slide_number}")

            # Extract text from shapes (including text boxes)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))

            full_text.append("\n".join(slide_text))

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        raise
