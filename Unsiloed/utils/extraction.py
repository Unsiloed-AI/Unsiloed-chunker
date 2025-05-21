import PyPDF2
import docx
from typing import Iterator, Union, BinaryIO
from pptx import Presentation
import os
import tempfile
import concurrent.futures
from urllib.parse import urlparse
import requests
import io
import logging
from .pdf_extraction import extract_text_streaming_robust

logger = logging.getLogger(__name__)

def is_url(path: str) -> bool:
    """Check if a path is a URL."""
    try:
        result = urlparse(path)
        return all([result.scheme, result.netloc])
    except:
        return False

def get_file_object(file_path: str) -> Union[BinaryIO, str]:
    """Get file object from path or URL."""
    if is_url(file_path):
        response = requests.get(file_path, stream=True)
        response.raise_for_status()
        return io.BytesIO(response.content)
    else:
        return file_path

def extract_text_streaming_pdf(file_path: Union[str, BinaryIO]) -> Iterator[dict]:
    """Extract text from PDF in a streaming fashion, yielding one page at a time.
    Optimized for performance with large documents and uses the best available extraction method."""
    
    # Try to import and use PyMuPDF for faster extraction if available
    try:
        from Unsiloed.utils.pdf_extraction import extract_text_streaming_pymupdf, PYMUPDF_AVAILABLE
        if PYMUPDF_AVAILABLE:
            yield from extract_text_streaming_pymupdf(file_path)
            return
    except Exception as e:
        logger.debug(f"PyMuPDF streaming not available: {str(e)}")
    
    # Fall back to robust extraction method if PyMuPDF is not available
    from Unsiloed.utils.pdf_extraction import extract_text_streaming_robust
    yield from extract_text_streaming_robust(file_path)

def extract_text_streaming_docx(file_path: Union[str, BinaryIO]) -> Iterator[str]:
    """Extract text from DOCX in a streaming fashion, yielding paragraphs in batches.
    Optimized for performance with better memory management and parallel processing for large documents."""
    file_obj = get_file_object(file_path) if isinstance(file_path, str) else file_path
    
    # If file_obj is a BytesIO or similar, we need to save it temporarily
    temp_file = None
    if not isinstance(file_obj, str):
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
        # Use a buffered approach to write large files
        if hasattr(file_obj, 'read'):
            # Stream the content in chunks to avoid loading the whole file in memory
            chunk_size = 2 * 1024 * 1024  # 2MB chunks for better performance
            with open(temp_file.name, 'wb') as f:
                while True:
                    chunk = file_obj.read(chunk_size)
                    if not chunk:
                        break
                    f.write(chunk)
        else:
            temp_file.write(file_obj.getvalue())
        temp_file.close()
        file_obj = temp_file.name
    
    try:
        doc = docx.Document(file_obj)
        paragraph_count = len(doc.paragraphs)
        
        # For very large documents, process paragraphs in batches to avoid memory issues
        large_document = paragraph_count > 1000
        
        # Use StringIO for more efficient string operations when concatenating text
        current_batch = io.StringIO()
        current_size = 0
        batch_size_limit = 8192 if large_document else 5000  # Larger batches for large docs
        batch_index = 0
        batch_para_start = 0
        
        for para_index, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue
                
            # Append to StringIO buffer instead of list for better performance
            current_batch.write(text)
            current_batch.write("\n")
            current_size += len(text) + 1  # +1 for the newline
            
            if current_size >= batch_size_limit:
                text_content = current_batch.getvalue()
                yield {
                    "text": text_content,
                    "metadata": {
                        "batch_index": batch_index,
                        "paragraph_range": f"{batch_para_start}-{para_index}",
                        "paragraphs": para_index - batch_para_start + 1,
                    }
                }
                # Reset buffer
                current_batch = io.StringIO()
                current_size = 0
                batch_index += 1
                batch_para_start = para_index + 1
        
        # Yield any remaining content
        if current_size > 0:
            text_content = current_batch.getvalue()
            yield {
                "text": text_content,
                "metadata": {
                    "batch_index": batch_index,
                    "paragraph_range": f"{batch_para_start}-{paragraph_count - 1}",
                    "paragraphs": paragraph_count - batch_para_start,
                }
            }
    finally:
        # Clean up resources
        if 'current_batch' in locals() and hasattr(current_batch, 'close'):
            current_batch.close()
            
        # Clean up temp file if created
        if temp_file and os.path.exists(file_obj):
            os.unlink(file_obj)

def extract_text_streaming_pptx(file_path: Union[str, BinaryIO]) -> Iterator[str]:
    """Extract text from PPTX in a streaming fashion, yielding one slide at a time."""
    file_obj = get_file_object(file_path) if isinstance(file_path, str) else file_path
    
    # If file_obj is a BytesIO or similar, we need to save it temporarily
    temp_file = None
    if not isinstance(file_obj, str):
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        temp_file.write(file_obj.read() if hasattr(file_obj, 'read') else file_obj.getvalue())
        temp_file.close()
        file_obj = temp_file.name
    
    try:
        presentation = Presentation(file_obj)
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                
                # Extract text from tables if present
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            if slide_text:
                yield {
                    "text": "\n".join(slide_text),
                    "metadata": {"slide": slide_num + 1}
                }
            else:
                yield {
                    "text": "",
                    "metadata": {"slide": slide_num + 1, "empty": True}
                }
    finally:
        # Clean up temp file if created
        if temp_file and os.path.exists(file_obj):
            os.unlink(file_obj)

def extract_text_streaming(file_path: str) -> Iterator[dict]:
    """
    Extract text from document in a streaming fashion, yielding content incrementally.
    Optimized for better performance and error handling.
    """
    # Use quick path detection with split once
    lower_path = file_path.lower()
    file_extension = os.path.splitext(lower_path)[1]

    # Fast dispatch based on extension
    if file_extension == '.pdf' or lower_path.endswith('.pdf'):
        yield from extract_text_streaming_pdf(file_path)
    elif file_extension == '.docx' or lower_path.endswith('.docx'):
        yield from extract_text_streaming_docx(file_path)
    elif file_extension == '.pptx' or lower_path.endswith('.pptx'):
        yield from extract_text_streaming_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_path}")
