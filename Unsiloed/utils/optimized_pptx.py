"""
Optimized PPTX processing module with memory-efficient operations.

This module provides advanced capabilities for handling PPTX files:
1. Streaming extraction for large presentations
2. Efficient text and image extraction
3. Structure preservation
4. Memory optimized processing
"""
import os
import io
import gc
import logging
import tempfile
from typing import Iterator, Dict, Union, BinaryIO, List, Any, Optional, Tuple, Generator
import concurrent.futures

from Unsiloed.utils.document_cache import document_cache
from Unsiloed.utils.memory_profiling import MemoryProfiler

logger = logging.getLogger(__name__)

# Try to import specialized presentation libraries
try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False
    logger.warning("python-pptx not available, will use alternative extraction methods")


def extract_text_streaming_pptx(
    file_path: Union[str, BinaryIO], 
    batch_size: int = 5
) -> Generator[Dict[str, Union[str, int]], None, None]:
    """
    Stream text extraction from PPTX with optimized memory usage.
    
    Args:
        file_path: Path to the PPTX file or file-like object
        batch_size: Number of slides to process in each batch
        
    Yields:
        Dictionary with slide number and extracted text
        
    Raises:
        PptxExtractionError: If there's an error during PPTX extraction
        DependencyError: If python-pptx is not available
    """
    # Try to import custom exceptions if available
    try:
        from Unsiloed.utils.exceptions import PptxExtractionError, DependencyError
        CUSTOM_EXCEPTIONS_AVAILABLE = True
    except ImportError:
        CUSTOM_EXCEPTIONS_AVAILABLE = False
    
    # Check if python-pptx is available
    if not PYTHON_PPTX_AVAILABLE:
        error_msg = "python-pptx is required for PPTX extraction"
        if CUSTOM_EXCEPTIONS_AVAILABLE:
            raise DependencyError(error_msg, dependency="python-pptx")
        else:
            raise RuntimeError(error_msg)
    
    # Use a memory profiler to track usage during streaming extraction
    profiler = MemoryProfiler(f"pptx_stream_extract_{os.path.basename(str(file_path))}")
    profiler.start()
    
    try:
        # Handle file-like objects by writing to temporary file if needed
        if isinstance(file_path, io.BytesIO):
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                temp_path = temp_file.name
                file_path.seek(0)
                temp_file.write(file_path.read())
            
            try:
                # Create presentation object
                try:
                    presentation = Presentation(temp_path)
                except Exception as e:
                    error_msg = f"Failed to open PPTX file from bytes: {str(e)}"
                    logger.error(error_msg)
                    if CUSTOM_EXCEPTIONS_AVAILABLE:
                        raise PptxExtractionError(error_msg, details={"error": str(e)}) from e
                    else:
                        raise RuntimeError(error_msg) from e
                
                # Get total slide count
                total_slides = len(presentation.slides)
                
                # Process slides in batches
                for i in range(0, total_slides, batch_size):
                    batch_end = min(i + batch_size, total_slides)
                    
                    for slide_num in range(i, batch_end):
                        slide = presentation.slides[slide_num]
                        try:
                            slide_text = extract_text_from_slide(slide)
                            
                            # Skip empty slides
                            if slide_text.strip():
                                yield {"slide": slide_num + 1, "text": slide_text}
                        except Exception as e:
                            # Log the error but continue processing other slides
                            error_msg = f"Error extracting text from slide {slide_num+1}: {str(e)}"
                            logger.warning(error_msg)
                    
                    # Explicit garbage collection after each batch
                    gc.collect()
            finally:
                # Clean up temporary file
                try:
                    os.unlink(temp_path)
                except Exception as e:
                    logger.debug(f"Failed to clean up temporary file: {str(e)}")
        else:
            # Direct file path processing
            try:
                presentation = Presentation(file_path)
            except Exception as e:
                error_msg = f"Failed to open PPTX file: {str(e)}"
                logger.error(error_msg)
                if CUSTOM_EXCEPTIONS_AVAILABLE:
                    raise PptxExtractionError(error_msg, file_path=str(file_path), details={"error": str(e)}) from e
                else:
                    raise RuntimeError(error_msg) from e
            
            # Get total slide count
            total_slides = len(presentation.slides)
            
            # Process slides in batches
            for i in range(0, total_slides, batch_size):
                batch_end = min(i + batch_size, total_slides)
                
                for slide_num in range(i, batch_end):
                    try:
                        slide = presentation.slides[slide_num]
                        slide_text = extract_text_from_slide(slide)
                        
                        # Skip empty slides
                        if slide_text.strip():
                            yield {"slide": slide_num + 1, "text": slide_text}
                    except Exception as e:
                        # Log the error but continue processing other slides
                        error_msg = f"Error extracting text from slide {slide_num+1}: {str(e)}"
                        logger.warning(error_msg)
                        # Yield placeholder for failed slide to maintain slide numbering
                        yield {"slide": slide_num + 1, "text": f"[Error extracting slide content: {str(e)}]"}
                
                # Explicit garbage collection after each batch
                gc.collect()
    except Exception as e:
        error_msg = f"Error during PPTX extraction: {str(e)}"
        logger.error(error_msg)
        if CUSTOM_EXCEPTIONS_AVAILABLE:
            raise PptxExtractionError(error_msg, file_path=str(file_path), details={"error": str(e)}) from e
        else:
            raise RuntimeError(error_msg) from e
    finally:
        # Stop profiling and log memory usage
        memory_stats = profiler.stop()
        logger.debug(f"Memory usage during PPTX extraction: {memory_stats}")


def extract_text_from_slide(slide):
    """
    Extract text from a PowerPoint slide.
    
    Args:
        slide: Slide object from python-pptx
        
    Returns:
        Extracted text as string
    """
    text_parts = []
    
    # Extract text from shapes
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text_parts.append(shape.text)
        
        # Handle tables
        if shape.has_table:
            for row in shape.table.rows:
                row_texts = []
                for cell in row.cells:
                    row_texts.append(cell.text)
                text_parts.append(" | ".join(row_texts))
    
    return "\n".join(text_parts)


def extract_pptx_with_structure(
    file_path: Union[str, BinaryIO]
) -> Dict[str, Any]:
    """
    Extract text from PPTX with structure preservation and memory optimization.
    
    Args:
        file_path: Path to the PPTX file or file-like object
        
    Returns:
        Dictionary with extracted structured content
        
    Raises:
        PptxExtractionError: If there's an error during PPTX extraction
        DependencyError: If python-pptx is not available
    """
    # Try to import custom exceptions if available
    try:
        from Unsiloed.utils.exceptions import PptxExtractionError, DependencyError, CacheError
        CUSTOM_EXCEPTIONS_AVAILABLE = True
    except ImportError:
        CUSTOM_EXCEPTIONS_AVAILABLE = False
    
    # Check if python-pptx is available
    if not PYTHON_PPTX_AVAILABLE:
        error_msg = "python-pptx is required for structured PPTX extraction"
        if CUSTOM_EXCEPTIONS_AVAILABLE:
            raise DependencyError(error_msg, dependency="python-pptx")
        else:
            raise RuntimeError(error_msg)
    
    # Check cache first using the proper pattern
    try:
        cached_result = document_cache.get(file_path, "pptx_structured")
        if cached_result:
            return cached_result
    except Exception as e:
        logger.debug(f"Cache lookup failed: {str(e)}")
    
    # Profile memory usage
    with MemoryProfiler(f"pptx_structured_{os.path.basename(str(file_path))}") as profiler:
        result = {
            "text": "",
            "slides": [],
            "metadata": {
                "total_slides": 0
            }
        }
        
        try:
            # Handle file-like objects by writing to temporary file if needed
            if isinstance(file_path, io.BytesIO):
                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                    temp_path = temp_file.name
                    file_path.seek(0)
                    temp_file.write(file_path.read())
                
                try:
                    # Process the presentation
                    try:
                        presentation = Presentation(temp_path)
                        result = extract_presentation_structure(presentation)
                    except Exception as e:
                        error_msg = f"Failed to process PPTX from bytes: {str(e)}"
                        logger.error(error_msg)
                        if CUSTOM_EXCEPTIONS_AVAILABLE:
                            raise PptxExtractionError(error_msg, details={"error": str(e)}) from e
                        else:
                            raise RuntimeError(error_msg) from e
                finally:
                    # Clean up temporary file
                    try:
                        os.unlink(temp_path)
                    except Exception as e:
                        logger.debug(f"Failed to clean up temporary file: {str(e)}")
            else:
                # Direct file path processing
                try:
                    presentation = Presentation(file_path)
                    result = extract_presentation_structure(presentation)
                except Exception as e:
                    error_msg = f"Failed to process PPTX file: {str(e)}"
                    logger.error(error_msg)
                    if CUSTOM_EXCEPTIONS_AVAILABLE:
                        raise PptxExtractionError(error_msg, file_path=str(file_path), details={"error": str(e)}) from e
                    else:
                        raise RuntimeError(error_msg) from e
            
            # Store in cache for future use
            try:
                document_cache.set(file_path, "pptx_structured", result)
            except Exception as e:
                logger.debug(f"Failed to cache result: {str(e)}")
                
            return result
        except Exception as e:
            if not isinstance(e, (PptxExtractionError, DependencyError)) and CUSTOM_EXCEPTIONS_AVAILABLE:
                error_msg = f"Error during PPTX extraction: {str(e)}"
                logger.error(error_msg)
                raise PptxExtractionError(error_msg, file_path=str(file_path), details={"error": str(e)}) from e
            raise
    
    return result


def extract_presentation_structure(presentation) -> Dict[str, Any]:
    """
    Extract structured data from a PowerPoint presentation.
    
    Args:
        presentation: Presentation object from python-pptx
        
    Returns:
        Dictionary with extracted structured content
    """
    result = {
        "text": "",
        "slides": [],
        "metadata": {
            "total_slides": len(presentation.slides)
        }
    }
    
    all_text = []
    slide_count = len(presentation.slides)
    
    # Process slides in batches to manage memory
    batch_size = 5
    for i in range(0, slide_count, batch_size):
        batch_end = min(i + batch_size, slide_count)
        
        for slide_num in range(i, batch_end):
            slide = presentation.slides[slide_num]
            
            # Extract all text from this slide
            slide_text = extract_text_from_slide(slide)
            
            # Extract shapes and their text
            shapes = []
            tables = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                # Handle text shapes
                if hasattr(shape, "text") and shape.text:
                    shape_data = {
                        "id": shape_idx,
                        "type": "text",
                        "text": shape.text
                    }
                    
                    # Try to determine if it's a title
                    if shape_idx == 0 or (hasattr(shape, "is_title") and shape.is_title):
                        shape_data["is_title"] = True
                    
                    shapes.append(shape_data)
                
                # Handle tables
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text)
                        table_data.append(row_data)
                    
                    tables.append({
                        "id": len(tables),
                        "data": table_data
                    })
            
            # Add slide data to result
            slide_data = {
                "number": slide_num + 1,
                "text": slide_text,
                "shapes": shapes,
                "tables": tables
            }
            
            result["slides"].append(slide_data)
            all_text.append(slide_text)
        
        # Explicit garbage collection after each batch
        gc.collect()
    
    # Combine all text
    result["text"] = "\n\n".join(all_text)
    
    return result


def extract_pptx_with_images(
    file_path: Union[str, BinaryIO]
) -> Dict[str, Any]:
    """
    Extract both text and images from a PPTX document with memory optimization.
    
    Args:
        file_path: Path to the PPTX file or file-like object
        
    Returns:
        Dictionary with extracted text and image data
    """
    # Check if python-pptx is available
    if not PYTHON_PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is required for image extraction")
    
    # Check cache first using the proper pattern
    try:
        cached_result = document_cache.get(file_path, "pptx_with_images")
        if cached_result:
            return cached_result
    except Exception as e:
        logger.debug(f"Cache lookup failed: {str(e)}")
    
    # Profile memory usage
    with MemoryProfiler(f"pptx_images_{os.path.basename(str(file_path))}") as profiler:
        # Get structured data first
        result = extract_pptx_with_structure(file_path)
        
        # Add images field if not present
        if "images" not in result:
            result["images"] = []
        
        # Create temporary directory for images
        with tempfile.TemporaryDirectory() as temp_dir:
            # Handle file-like objects by writing to temporary file if needed
            if isinstance(file_path, io.BytesIO):
                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                    temp_path = temp_file.name
                    file_path.seek(0)
                    temp_file.write(file_path.read())
                file_to_use = temp_path
            else:
                file_to_use = file_path
            
            try:
                # Process the presentation for images
                presentation = Presentation(file_to_use)
                
                # Process slides in batches
                slide_count = len(presentation.slides)
                batch_size = 3  # Smaller batch size for image extraction
                
                for i in range(0, slide_count, batch_size):
                    batch_end = min(i + batch_size, slide_count)
                    
                    for slide_num in range(i, batch_end):
                        slide = presentation.slides[slide_num]
                        
                        # Extract images from shapes
                        for shape_idx, shape in enumerate(slide.shapes):
                            # Check if shape has an image
                            if hasattr(shape, "image") and shape.image:
                                try:
                                    # Get image data
                                    image = shape.image
                                    
                                    # Create image info
                                    image_info = {
                                        "slide": slide_num + 1,
                                        "shape_id": shape_idx,
                                        "content_type": image.content_type,
                                        "size": len(image.blob)
                                    }
                                    
                                    # Try to get dimensions if PIL is available
                                    try:
                                        from PIL import Image
                                        image_stream = io.BytesIO(image.blob)
                                        with Image.open(image_stream) as img:
                                            image_info["width"] = img.width
                                            image_info["height"] = img.height
                                    except:
                                        pass
                                    
                                    # Add to results
                                    result["images"].append(image_info)
                                except Exception as e:
                                    logger.warning(f"Error extracting image: {str(e)}")
                    
                    # Explicit garbage collection after each batch
                    gc.collect()
            finally:
                # Clean up temporary file if created
                if isinstance(file_path, io.BytesIO):
                    try:
                        os.unlink(file_to_use)
                    except:
                        pass
    
    # Store in cache for future use
    try:
        document_cache.set(file_path, "pptx_with_images", result)
    except Exception as e:
        logger.debug(f"Failed to cache result: {str(e)}")
    
    return result
