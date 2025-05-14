"""
Model-agnostic utilities for OCR and text processing.
This module provides functions for working with different model providers.
"""

import os
import base64
import json
from typing import List, Dict, Any, Optional
import logging
import concurrent.futures
import PyPDF2
from dotenv import load_dotenv
import numpy as np
import cv2

from Unsiloed.models import get_provider
from Unsiloed.config import config

load_dotenv()

logger = logging.getLogger(__name__)

# Configure logging
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)


def encode_image_to_base64(image_path):
    """
    Encode an image to base64.

    Args:
        image_path: Path to the image file or numpy array

    Returns:
        Base64 encoded string of the image
    """
    logger.debug("Encoding image to base64")

    # Handle numpy array (from CV2)
    if isinstance(image_path, np.ndarray):
        success, buffer = cv2.imencode(".jpg", image_path)
        if not success:
            raise ValueError("Failed to encode image")
        return base64.b64encode(buffer).decode("utf-8")

    # Handle file path
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def create_extraction_prompt(schema: Dict[str, Any], page_count: int) -> str:
    """
    Create a prompt instructing the model how to extract data according to the schema.

    Args:
        schema: JSON schema defining the structure
        page_count: Number of pages in the document

    Returns:
        Prompt string for the model
    """
    # Create a compact JSON representation of the schema
    schema_str = json.dumps(schema, indent=2)

    prompt = f"""
    You are an expert at extracting structured data from documents.

    I have a document with {page_count} pages that has been converted to images. I need you to extract specific information from these images according to the following JSON schema:

    {schema_str}

    Please follow these instructions carefully:

    1. Examine all {page_count} images thoroughly to find the requested information.
    2. Extract the exact text from the document that matches each field in the schema.
    3. If you cannot find information for a specific field in any of the pages, return an empty string or null value for that field.
    4. For array fields, include all instances found throughout the document.
    5. Maintain the structure defined in the schema exactly.
    6. Return only the extracted data as a valid JSON object, matching the structure of the schema.
    7. Do not add any explanatory text or notes outside the JSON structure.
    8. Be precise and accurate in your extraction.
    9. If text is unclear or ambiguous, make your best guess based on context.
    10. For dates, numbers, and other formatted data, maintain the format as shown in the document.
    11. IMPORTANT: Your response MUST be a valid JSON object that exactly matches the structure of the provided schema.
    12. IMPORTANT: Do not include any explanations, just return the JSON object.

    Your response should be a valid JSON object containing only the extracted data.
    """

    return prompt


def semantic_chunk_with_structured_output(text: str, provider_name: Optional[str] = None) -> List[Dict[str, Any]]:
    """
    Use the selected model provider to create semantic chunks from text.

    Args:
        text: The text to chunk
        provider_name: Name of the model provider to use (default: from config)

    Returns:
        List of chunks with metadata
    """
    # If text is too long, split it first using a simpler method
    # and then process each part in parallel
    if len(text) > 25000:
        logger.info(
            "Text too long for direct semantic chunking, applying parallel processing"
        )
        return process_long_text_semantically(text, provider_name)

    try:
        # Get the model provider
        provider_config = config.get_provider_config(provider_name)
        model_provider = get_provider(provider_name, **provider_config)

        # Create a system prompt for the model
        system_prompt = "You are an expert at analyzing and dividing text into meaningful semantic chunks. Your output should be valid JSON."

        # Define the output schema
        output_schema = {
            "chunks": [
                {
                    "text": "the text of the chunk",
                    "title": "a descriptive title for this chunk",
                    "position": "beginning/middle/end"
                }
            ]
        }

        # Create the user prompt
        user_prompt = f"""Please analyze the following text and divide it into logical semantic chunks.
        Each chunk should represent a cohesive unit of information or a distinct section.

        Text to chunk:

        {text}"""

        # Generate structured output using the model provider
        result = model_provider.generate_structured_output(
            prompt=user_prompt,
            system_prompt=system_prompt,
            output_schema=output_schema,
            temperature=0.1
        )

        # Convert the response to our standard chunk format
        chunks = []
        current_position = 0

        for i, chunk_data in enumerate(result.get("chunks", [])):
            chunk_text = chunk_data.get("text", "")
            # Find the chunk in the original text to get accurate character positions
            start_position = text.find(chunk_text, current_position)
            if start_position == -1:
                # If exact match not found, use approximate position
                start_position = current_position

            end_position = start_position + len(chunk_text)

            chunks.append(
                {
                    "text": chunk_text,
                    "metadata": {
                        "title": chunk_data.get("title", f"Chunk {i + 1}"),
                        "position": chunk_data.get("position", "unknown"),
                        "start_char": start_position,
                        "end_char": end_position,
                        "strategy": "semantic",
                    },
                }
            )

            current_position = end_position

        return chunks

    except Exception as e:
        logger.error(f"Error in semantic chunking with model provider: {str(e)}")
        # Fall back to paragraph chunking if semantic chunking fails
        logger.info("Falling back to paragraph chunking")
        # We'll just do basic paragraph chunking here
        paragraphs = text.split("\n\n")
        paragraphs = [p.strip() for p in paragraphs if p.strip()]

        chunks = []
        current_position = 0

        for i, paragraph in enumerate(paragraphs):
            start_position = text.find(paragraph, current_position)
            if start_position == -1:
                start_position = current_position

            end_position = start_position + len(paragraph)

            chunks.append(
                {
                    "text": paragraph,
                    "metadata": {
                        "title": f"Paragraph {i + 1}",
                        "position": "unknown",
                        "start_char": start_position,
                        "end_char": end_position,
                        "strategy": "paragraph",  # Fall back strategy
                    },
                }
            )

            current_position = end_position

        return chunks


def process_long_text_semantically(text: str, provider_name: Optional[str] = None) -> List[Dict[str, Any]]:
    """
    Process a long text by breaking it into smaller pieces and chunking each piece semantically.
    Uses parallel processing for better performance.

    Args:
        text: The long text to process
        provider_name: Name of the model provider to use (default: from config)

    Returns:
        List of semantic chunks
    """
    # Create chunks of 25000 characters with 500 character overlap
    text_chunks = []
    chunk_size = 25000
    overlap = 500
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(start + chunk_size, text_length)
        text_chunks.append(text[start:end])
        start = end - overlap if end < text_length else text_length

    # Process each chunk in parallel
    all_semantic_chunks = []
    with concurrent.futures.ThreadPoolExecutor() as executor:
        # Define a worker function
        def process_chunk(chunk_text):
            try:
                # Process this chunk with the selected model provider
                sub_chunks = semantic_chunk_with_structured_output(chunk_text, provider_name)
                return sub_chunks
            except Exception as e:
                logger.error(
                    f"Error processing semantic subchunk: {str(e)}"
                )
                return []

        # Submit all tasks and gather results
        futures = [executor.submit(process_chunk, chunk) for chunk in text_chunks]
        for future in concurrent.futures.as_completed(futures):
            all_semantic_chunks.extend(future.result())

    return all_semantic_chunks


# Document text extraction functions (moved from openai.py)

def extract_text_from_pdf(pdf_path: str) -> str:
    """
    Extract text from a PDF file with optimized performance.
    Uses parallel processing for multi-page PDFs.

    Args:
        pdf_path: Path to the PDF file

    Returns:
        Extracted text from the PDF
    """
    try:
        with open(pdf_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)

            # Function to extract text from a page
            def extract_page_text(page_idx):
                try:
                    page = reader.pages[page_idx]
                    text = page.extract_text() or ""
                    return text
                except Exception as e:
                    logger.warning(
                        f"Error extracting text from page {page_idx}: {str(e)}"
                    )
                    return ""

            # For small PDFs, sequential processing is faster
            if len(reader.pages) <= 5:
                all_text = ""
                for i in range(len(reader.pages)):
                    all_text += extract_page_text(i) + "\n\n"
            else:
                # Process pages in parallel for larger PDFs
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    results = list(
                        executor.map(extract_page_text, range(len(reader.pages)))
                    )
                all_text = "\n\n".join(results)

        return all_text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        raise


def extract_text_from_docx(docx_path: str) -> str:
    """
    Extract text from a DOCX file.

    Args:
        docx_path: Path to the DOCX file

    Returns:
        Extracted text from the DOCX
    """
    try:
        import docx  # python-docx package

        doc = docx.Document(docx_path)
        full_text = []

        # Extract text from paragraphs
        for para in doc.paragraphs:
            full_text.append(para.text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(cell.text)

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        raise


def extract_text_from_pptx(pptx_path: str) -> str:
    """
    Extract text from a PPTX file.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        Extracted text from the PPTX
    """
    try:
        from pptx import Presentation  # python-pptx package

        presentation = Presentation(pptx_path)
        full_text = []

        # Loop through slides
        for slide_number, slide in enumerate(presentation.slides, 1):
            slide_text = []
            slide_text.append(f"Slide {slide_number}")

            # Extract text from shapes (including text boxes)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))

            full_text.append("\n".join(slide_text))

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        raise


def extract_text_from_doc(doc_path: str) -> str:
    """
    Extract text from a DOC file (older Microsoft Word format).

    Args:
        doc_path: Path to the DOC file

    Returns:
        Extracted text from the DOC
    """
    try:
        import docx2txt  # docx2txt package

        # Extract text including images (images will be ignored)
        text = docx2txt.process(doc_path)
        return text
    except Exception as e:
        logger.error(f"Error extracting text from DOC: {str(e)}")
        raise


def extract_text_from_xlsx(xlsx_path: str) -> str:
    """
    Extract text from an XLSX file (Excel).

    Args:
        xlsx_path: Path to the XLSX file

    Returns:
        Extracted text from the XLSX
    """
    try:
        import openpyxl  # openpyxl package

        workbook = openpyxl.load_workbook(xlsx_path, data_only=True)
        full_text = []

        # Process each worksheet
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = [f"Sheet: {sheet_name}"]

            # Get the used range
            rows = list(sheet.rows)
            if not rows:
                continue

            # Process each row
            for row in rows:
                row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                if any(row_values):  # Skip empty rows
                    sheet_text.append(" | ".join(row_values))

            full_text.append("\n".join(sheet_text))

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from XLSX: {str(e)}")
        raise


def extract_text_from_xls(xls_path: str) -> str:
    """
    Extract text from an XLS file (older Excel format).

    Args:
        xls_path: Path to the XLS file

    Returns:
        Extracted text from the XLS
    """
    try:
        import xlrd  # xlrd package

        workbook = xlrd.open_workbook(xls_path)
        full_text = []

        # Process each worksheet
        for sheet_index in range(workbook.nsheets):
            sheet = workbook.sheet_by_index(sheet_index)
            sheet_text = [f"Sheet: {sheet.name}"]

            # Process each row
            for row_idx in range(sheet.nrows):
                row_values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                if any(value.strip() for value in row_values):  # Skip empty rows
                    sheet_text.append(" | ".join(row_values))

            full_text.append("\n".join(sheet_text))

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from XLS: {str(e)}")
        raise


def extract_text_from_odt(odt_path: str) -> str:
    """
    Extract text from an ODT file (OpenDocument Text).

    Args:
        odt_path: Path to the ODT file

    Returns:
        Extracted text from the ODT
    """
    try:
        from odf import text, teletype
        from odf.opendocument import load

        textdoc = load(odt_path)
        paragraphs = textdoc.getElementsByType(text.P)

        # Extract text from paragraphs
        full_text = [teletype.extractText(paragraph) for paragraph in paragraphs]

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from ODT: {str(e)}")
        raise


def extract_text_from_ods(ods_path: str) -> str:
    """
    Extract text from an ODS file (OpenDocument Spreadsheet).

    Args:
        ods_path: Path to the ODS file

    Returns:
        Extracted text from the ODS
    """
    try:
        from odf import text, teletype
        from odf.opendocument import load
        from odf.table import Table, TableRow, TableCell

        spreadsheet = load(ods_path)
        full_text = []

        # Get all tables
        tables = spreadsheet.getElementsByType(Table)

        for table_index, table in enumerate(tables):
            table_text = [f"Sheet {table_index + 1}: {table.getAttribute('name') or ''}"]

            # Get all rows in the table
            rows = table.getElementsByType(TableRow)

            for row in rows:
                # Get all cells in the row
                cells = row.getElementsByType(TableCell)
                row_values = []

                for cell in cells:
                    # Extract text from the cell
                    paragraphs = cell.getElementsByType(text.P)
                    cell_text = " ".join([teletype.extractText(p) for p in paragraphs])
                    row_values.append(cell_text)

                if any(row_values):  # Skip empty rows
                    table_text.append(" | ".join(row_values))

            full_text.append("\n".join(table_text))

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from ODS: {str(e)}")
        raise


def extract_text_from_odp(odp_path: str) -> str:
    """
    Extract text from an ODP file (OpenDocument Presentation).

    Args:
        odp_path: Path to the ODP file

    Returns:
        Extracted text from the ODP
    """
    try:
        from odf import text, teletype
        from odf.opendocument import load
        from odf.draw import Page

        presentation = load(odp_path)
        full_text = []

        # Get all pages (slides)
        pages = presentation.getElementsByType(Page)

        for page_index, page in enumerate(pages):
            page_text = [f"Slide {page_index + 1}: {page.getAttribute('name') or ''}"]

            # Get all text elements in the page
            paragraphs = page.getElementsByType(text.P)

            for paragraph in paragraphs:
                paragraph_text = teletype.extractText(paragraph)
                if paragraph_text.strip():
                    page_text.append(paragraph_text)

            full_text.append("\n".join(page_text))

        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extracting text from ODP: {str(e)}")
        raise


def extract_text_from_txt(txt_path: str) -> str:
    """
    Extract text from a TXT file (plain text).

    Args:
        txt_path: Path to the TXT file

    Returns:
        Extracted text from the TXT
    """
    try:
        with open(txt_path, 'r', encoding='utf-8', errors='replace') as file:
            return file.read()
    except Exception as e:
        logger.error(f"Error extracting text from TXT: {str(e)}")
        raise


def extract_text_from_rtf(rtf_path: str) -> str:
    """
    Extract text from an RTF file (Rich Text Format).

    Args:
        rtf_path: Path to the RTF file

    Returns:
        Extracted text from the RTF
    """
    try:
        from striprtf.striprtf import rtf_to_text

        with open(rtf_path, 'r', encoding='utf-8', errors='replace') as file:
            rtf_text = file.read()

        # Convert RTF to plain text
        plain_text = rtf_to_text(rtf_text)
        return plain_text
    except Exception as e:
        logger.error(f"Error extracting text from RTF: {str(e)}")
        raise


def extract_text_from_epub(epub_path: str) -> str:
    """
    Extract text from an EPUB file (e-book).

    Args:
        epub_path: Path to the EPUB file

    Returns:
        Extracted text from the EPUB
    """
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup

        # Function to extract text from HTML content
        def chapter_to_text(chapter_content):
            soup = BeautifulSoup(chapter_content, 'html.parser')
            text = soup.get_text()
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            return text

        # Load the EPUB file
        book = epub.read_epub(epub_path)
        chapters = []

        # Extract content from each chapter
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                chapters.append(chapter_to_text(item.get_content()))

        return "\n\n".join(chapters)
    except Exception as e:
        logger.error(f"Error extracting text from EPUB: {str(e)}")
        raise
